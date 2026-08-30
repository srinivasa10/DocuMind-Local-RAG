from __future__ import annotations

import csv
import json
import logging
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class PageText:
    page: int | None
    text: str


def _ocr_page_image(image_bytes: bytes, mime_type: str = "image/png") -> str:
    """Extract text from an image or scanned document page using Gemini Vision OCR.

    Supports image mime types (image/png, image/jpeg, image/webp) as well as
    application/pdf for direct PDF document understanding.
    """
    from app.config import get_settings

    settings = get_settings()
    if not settings.gemini_api_key:
        return ""
    try:
        from google import genai
        from google.genai import types

        client = genai.Client(api_key=settings.gemini_api_key)
        response = client.models.generate_content(
            model=settings.gemini_model,
            contents=[
                types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                (
                    "Extract ALL text verbatim from this document. Include every section, heading, "
                    "paragraph, bullet point, table, contact info, skills, experience, education, "
                    "projects, certifications, and any other visible content. "
                    "Preserve the original structure as closely as possible."
                ),
            ],
        )
        return response.text or ""
    except Exception as exc:
        logger.warning("Vision OCR extraction failed: %s", exc)
        return ""


def _ocr_full_pdf(path: Path) -> list[PageText]:
    """Send the entire PDF file to Gemini as a PDF document for full-document OCR.

    Used as a last-resort fallback when neither PyMuPDF nor pypdf can extract text
    (e.g. scanned / image-only PDFs running without PyMuPDF installed).
    """
    from app.config import get_settings

    settings = get_settings()
    if not settings.gemini_api_key:
        logger.warning("Cannot perform PDF OCR: GEMINI_API_KEY not set")
        return [PageText(page=1, text="")]
    try:
        pdf_bytes = path.read_bytes()
        logger.info("Sending full PDF %s (%d bytes) to Gemini for document-level OCR...", path.name, len(pdf_bytes))
        text = _ocr_page_image(pdf_bytes, mime_type="application/pdf")
        if text.strip():
            return [PageText(page=1, text=text.strip())]
    except Exception as exc:
        logger.warning("Full PDF OCR failed for %s: %s", path.name, exc)
    return [PageText(page=1, text="")]


def _extract_pdf(path: Path) -> list[PageText]:
    """Extract text from PDF with a 3-tier fallback chain.

    Tier 1 — PyMuPDF native text extraction (fastest, most accurate for text PDFs).
    Tier 2 — PyMuPDF page-by-page image render + Gemini Vision OCR (scanned PDFs).
    Tier 3 — pypdf text extraction + per-image Gemini Vision OCR fallback.
    Tier 4 — Full document Gemini PDF understanding (last resort, no PyMuPDF needed).
    """
    pymupdf_available = False
    try:
        try:
            import pymupdf as fitz  # PyMuPDF >= 1.24 modern import
        except ImportError:
            import fitz  # type: ignore[no-redef]  # PyMuPDF legacy import
        pymupdf_available = True

        doc = fitz.open(str(path))
        pages: list[PageText] = []
        for index, page in enumerate(doc):
            # "text" mode preserves reading order and multi-column layouts
            text = page.get_text("text") or ""
            pages.append(PageText(page=index + 1, text=text.strip()))

        # --- Tier 1: embedded text found ---
        if any(p.text for p in pages):
            logger.info("PDF %s: extracted text via PyMuPDF native (%d pages)", path.name, len(pages))
            return pages

        # --- Tier 2: scanned PDF — render each page and OCR ---
        logger.info(
            "PDF %s: no embedded text; rendering %d pages for Vision OCR...",
            path.name, len(doc),
        )
        ocr_pages: list[PageText] = []
        for index, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            extracted = _ocr_page_image(img_bytes, mime_type="image/png")
            ocr_pages.append(PageText(page=index + 1, text=extracted.strip()))

        if any(p.text for p in ocr_pages):
            logger.info("PDF %s: extracted text via PyMuPDF + Vision OCR", path.name)
            return ocr_pages

    except Exception as exc:
        logger.warning(
            "PyMuPDF extraction failed for %s (%s); falling back to pypdf", path.name, exc
        )

    # --- Tier 3: pypdf text + per-image OCR ---
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = [
            PageText(page=index + 1, text=(page.extract_text() or "").strip())
            for index, page in enumerate(reader.pages)
        ]

        if any(p.text for p in pages):
            logger.info("PDF %s: extracted text via pypdf (%d pages)", path.name, len(pages))
            return pages

        # pypdf returned empty — try extracting embedded images and OCRing them
        logger.info("PDF %s: pypdf returned empty text; trying per-page image OCR...", path.name)
        ocr_pages = []
        for index, page in enumerate(reader.pages):
            page_blocks: list[str] = []
            for img in getattr(page, "images", []):
                try:
                    extracted = _ocr_page_image(img.data, mime_type="image/png")
                    if extracted.strip():
                        page_blocks.append(extracted.strip())
                except Exception:
                    pass
            ocr_pages.append(PageText(page=index + 1, text="\n\n".join(page_blocks)))

        if any(p.text for p in ocr_pages):
            logger.info("PDF %s: extracted text via pypdf image OCR", path.name)
            return ocr_pages

    except Exception as exc:
        logger.warning("pypdf extraction failed for %s (%s)", path.name, exc)

    # --- Tier 4: send entire PDF to Gemini for document-level understanding ---
    logger.info("PDF %s: all text methods exhausted; attempting full-document Gemini OCR...", path.name)
    result = _ocr_full_pdf(path)
    if any(p.text for p in result):
        logger.info("PDF %s: extracted text via full-document Gemini OCR", path.name)
        return result

    logger.error("PDF %s: all extraction methods failed", path.name)
    return [PageText(page=1, text="")]


def _extract_image(path: Path) -> list[PageText]:
    """Extract text from direct image files (.png, .jpg, .jpeg, .webp)."""
    mime_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
    }
    mime = mime_map.get(path.suffix.lower(), "image/png")
    image_bytes = path.read_bytes()
    text = _ocr_page_image(image_bytes, mime_type=mime)
    return [PageText(page=1, text=text)]


def _extract_docx(path: Path) -> list[PageText]:
    """Extract text from Microsoft Word .docx files including tables."""
    import docx

    doc = docx.Document(str(path))
    content_blocks: list[str] = []

    for p in doc.paragraphs:
        if p.text.strip():
            content_blocks.append(p.text.strip())

    for table in doc.tables:
        table_rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                table_rows.append(" | ".join(cells))
        if table_rows:
            content_blocks.append("\n".join(table_rows))

    full_text = "\n\n".join(content_blocks)
    return [PageText(page=1, text=full_text)]


def _extract_pptx(path: Path) -> list[PageText]:
    """Extract text from PowerPoint presentations slide-by-slide."""
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[PageText] = []

    for index, slide in enumerate(prs.slides):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_lines.append(line)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"Notes: {notes}")
        pages.append(PageText(page=index + 1, text="\n".join(slide_lines)))

    return pages


def _extract_excel(path: Path) -> list[PageText]:
    """Extract tabular text from Excel .xlsx / .xls workbooks."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    pages: list[PageText] = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        sheet = wb[sheet_name]
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            if any(val is not None for val in row):
                row_str = " | ".join(str(val) if val is not None else "" for val in row)
                rows.append(row_str)
        if rows:
            pages.append(PageText(page=sheet_idx + 1, text=f"Sheet: {sheet_name}\n" + "\n".join(rows)))

    return pages if pages else [PageText(page=1, text="")]


def _extract_csv_tsv(path: Path, is_tsv: bool = False) -> list[PageText]:
    """Extract tabular text from CSV / TSV files."""
    delimiter = "\t" if is_tsv else ","
    rows: list[str] = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(cell.strip() for cell in row))
    return [PageText(page=1, text="\n".join(rows))]


def _extract_json(path: Path) -> list[PageText]:
    """Extract structured text from JSON files."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        parsed = json.loads(raw)
        formatted = json.dumps(parsed, indent=2)
        return [PageText(page=1, text=formatted)]
    except Exception:
        return [PageText(page=1, text=raw)]


def _extract_plain(path: Path) -> list[PageText]:
    """Extract text from plain text, markdown, yaml, html, etc."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    return [PageText(page=1, text=raw)]


def extract_text(path: Path) -> list[PageText]:
    """Universal document text extractor supporting PDF, DOCX, XLSX, PPTX, CSV, JSON, and Text formats."""
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix in {".xlsx", ".xls"}:
        return _extract_excel(path)
    if suffix in {".csv", ".tsv"}:
        return _extract_csv_tsv(path, is_tsv=(suffix == ".tsv"))
    if suffix == ".json":
        return _extract_json(path)
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return _extract_image(path)
    if suffix in {".txt", ".md", ".markdown", ".yaml", ".yml", ".html", ".xml", ".log", ".rst"}:
        return _extract_plain(path)

    # Fallback to plain text read
    try:
        return _extract_plain(path)
    except Exception as exc:
        raise ValueError(f"Unsupported file format: {suffix}") from exc


